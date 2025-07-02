#!/usr/bin/env python3
import sys
import os
import argparse
import json

def pdf_to_word(pdf_file_path, docx_file_path):
    """
    Convert PDF to DOCX using pdf2docx library
    Based on working code from user's Gradio app
    """
    try:
        from pdf2docx import Converter
        print(f"Starting conversion of {pdf_file_path} to {docx_file_path}")
        
        # Create converter object
        cv = Converter(pdf_file_path)
        
        # Convert PDF to DOCX with multi-processing
        cv.convert(docx_file_path, multi_processing=True, start=0, end=None)
        
        # Close the converter
        cv.close()
        
        print(f"Successfully converted {pdf_file_path} to {docx_file_path}")
        return True
        
    except Exception as e:
        print(f"Error converting PDF to DOCX: {str(e)}")
        return False

def pdf_to_excel(pdf_file_path, excel_file_path, page_selection='all'):
    """
    Convert PDF to Excel using pdfplumber, pandas, and openpyxl
    Advanced table extraction with multiple refinements and page selection
    """
    try:
        import pdfplumber
        import pandas as pd
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        
        print(f"Starting conversion of {pdf_file_path} to {excel_file_path}")
        print(f"Page selection: {page_selection}")
        
        # Parse page selection
        selected_pages = parse_page_selection(page_selection)
        print(f"Selected pages: {selected_pages}")
        
        # Extract tables and text from PDF with advanced detection
        tables = []
        text_data = []
        
        with pdfplumber.open(pdf_file_path) as pdf:
            total_pages = len(pdf.pages)
            print(f"Total pages in PDF: {total_pages}")
            
            for page_num, page in enumerate(pdf.pages):
                current_page = page_num + 1
                
                # Skip pages not in selection
                if selected_pages and current_page not in selected_pages:
                    print(f"Skipping page {current_page} (not in selection)")
                    continue
                
                print(f"Processing page {current_page}")
                
                # Advanced table extraction with multiple strategies
                page_tables = []
                
                # Strategy 1: Extract tables with default settings
                try:
                    default_tables = page.extract_tables()
                    if default_tables:
                        page_tables.extend(default_tables)
                        print(f"Strategy 1 found {len(default_tables)} tables")
                except Exception as e:
                    print(f"Strategy 1 failed: {str(e)}")
                
                # Strategy 2: Extract tables with custom settings for better detection
                try:
                    custom_tables = page.extract_tables({
                        'vertical_strategy': 'text',
                        'horizontal_strategy': 'text',
                        'intersection_x_tolerance': 10,
                        'intersection_y_tolerance': 10
                    })
                    if custom_tables:
                        page_tables.extend(custom_tables)
                        print(f"Strategy 2 found {len(custom_tables)} tables")
                except Exception as e:
                    print(f"Strategy 2 failed: {str(e)}")
                
                # Strategy 3: Extract tables with lines strategy (more robust)
                try:
                    lines_tables = page.extract_tables({
                        'vertical_strategy': 'lines',
                        'horizontal_strategy': 'lines'
                    })
                    if lines_tables:
                        page_tables.extend(lines_tables)
                        print(f"Strategy 3 found {len(lines_tables)} tables")
                except Exception as e:
                    print(f"Strategy 3 failed: {str(e)}")
                
                # Strategy 4: Extract tables with minimal settings
                try:
                    minimal_tables = page.extract_tables({
                        'vertical_strategy': 'text',
                        'horizontal_strategy': 'text'
                    })
                    if minimal_tables:
                        page_tables.extend(minimal_tables)
                        print(f"Strategy 4 found {len(minimal_tables)} tables")
                except Exception as e:
                    print(f"Strategy 4 failed: {str(e)}")
                
                # Remove duplicates and process tables
                unique_tables = []
                for table in page_tables:
                    if table and len(table) > 0:
                        # Clean and validate table
                        cleaned_table = clean_table_data(table)
                        if cleaned_table and len(cleaned_table) > 0:
                            # Check if this table is unique (not a duplicate)
                            if not is_duplicate_table(cleaned_table, unique_tables):
                                unique_tables.append(cleaned_table)
                
                # Add unique tables to main list
                for table_idx, table in enumerate(unique_tables):
                    tables.append({
                        'page': current_page,
                        'table_index': table_idx + 1,
                        'data': table,
                        'rows': len(table),
                        'cols': len(table[0]) if table else 0
                    })
                    print(f"Found table {table_idx + 1} on page {current_page} with {len(table)} rows and {len(table[0]) if table else 0} columns")
                
                # Extract text if no tables found
                text = page.extract_text()
                if text and text.strip():
                    text_data.append({
                        'Page': current_page,
                        'Content': text.strip()
                    })
        
        print(f"Found {len(tables)} unique tables and {len(text_data)} text sections")
        
        # Create Excel workbook
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Process tables with advanced formatting
        if tables:
            for table_idx, table_info in enumerate(tables):
                # Create sheet for each table with page info
                if table_idx == 0:
                    ws = wb.create_sheet(f"Page_{table_info['page']}_Table_1")
                else:
                    ws = wb.create_sheet(f"Page_{table_info['page']}_Table_{table_info['table_index']}")
                
                table_data = table_info['data']
                
                # Add table data with advanced formatting
                for row_idx, row in enumerate(table_data, start=1):
                    for col_idx, cell_value in enumerate(row, start=1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                        
                        # Apply cell formatting
                        if row_idx == 1:  # Header row
                            cell.font = Font(bold=True, size=12)
                            cell.alignment = Alignment(horizontal='center', vertical='center')
                        else:
                            cell.alignment = Alignment(horizontal='left', vertical='top')
                
                # Add borders to table
                add_table_borders(ws, len(table_data), len(table_data[0]) if table_data else 0)
                
                # Auto-adjust column widths with smart sizing
                adjust_column_widths(ws, table_data)
                
                print(f"Added table {table_idx + 1} to sheet {ws.title} with {len(table_data)} rows")
        
        # If no tables found, create user-friendly message sheet
        if not tables:
            ws = wb.create_sheet("No_Tables_Found")
            
            # Add clear user-friendly message
            ws['A1'] = "No Tables Found in Selected Pages"
            ws['A1'].font = Font(bold=True, size=14, color="FF0000")  # Red color for emphasis
            
            ws['A3'] = "What this means:"
            ws['A3'].font = Font(bold=True, size=12)
            
            ws['A4'] = f"• No tables found in the selected pages: {page_selection}"
            ws['A5'] = "• The content might be text-only or image-based"
            ws['A6'] = "• Tables might be embedded as images"
            ws['A7'] = "• The layout might not be recognized as tables"
            
            ws['A9'] = "Suggestions:"
            ws['A9'].font = Font(bold=True, size=12)
            
            ws['A10'] = "• Try selecting different pages"
            ws['A11'] = "• Try a PDF with clear table borders"
            ws['A12'] = "• Ensure tables have visible grid lines"
            ws['A13'] = "• Check if tables are not just text with spaces"
            ws['A14'] = "• Try using PDF to Word conversion instead"
            
            ws['A16'] = "Alternative:"
            ws['A16'].font = Font(bold=True, size=12)
            
            ws['A17'] = "If you need the text content, try PDF to Word conversion"
            ws['A18'] = "which can better handle text extraction and formatting."
            
            # Auto-adjust column width
            ws.column_dimensions['A'].width = 60
            
            print("No tables found - created user-friendly message sheet")
        
        # If no tables found but text exists, create text sheet
        if not tables and text_data:
            ws = wb.create_sheet("Text_Content")
            
            # Add headers
            ws['A1'] = "Page"
            ws['B1'] = "Content"
            
            # Style headers
            header_font = Font(bold=True)
            ws['A1'].font = header_font
            ws['B1'].font = header_font
            
            # Add text data
            for row_idx, item in enumerate(text_data, start=2):
                ws[f'A{row_idx}'] = item['Page']
                ws[f'B{row_idx}'] = item['Content']
            
            # Auto-adjust column widths
            ws.column_dimensions['A'].width = 8
            ws.column_dimensions['B'].width = 80
            
            print("Added text content to sheet")
        
        # If still no data, create empty sheet with message
        if not tables and not text_data:
            ws = wb.create_sheet("No_Data_Found")
            ws['A1'] = "No Extractable Data Found"
            ws['A1'].font = Font(bold=True, size=14, color="FF0000")
            
            ws['A3'] = "The PDF appears to be:"
            ws['A3'].font = Font(bold=True, size=12)
            
            ws['A4'] = "• Image-based (scanned document)"
            ws['A5'] = "• Password protected"
            ws['A6'] = "• Corrupted or unreadable"
            ws['A7'] = "• Empty or contains no text"
            
            ws['A9'] = "Recommendations:"
            ws['A9'].font = Font(bold=True, size=12)
            
            ws['A10'] = "• Use OCR software to extract text from images"
            ws['A11'] = "• Check if the PDF is password protected"
            ws['A12'] = "• Try a different PDF file"
            ws['A13'] = "• Ensure the PDF is not corrupted"
            
            ws.column_dimensions['A'].width = 50
            
            print("Created empty sheet with detailed message")
        
        # Save the workbook
        wb.save(excel_file_path)
        print(f"Successfully created Excel file: {excel_file_path}")
        return True
        
    except Exception as e:
        print(f"Error converting PDF to Excel: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def clean_table_data(table):
    """
    Clean and validate table data
    """
    if not table:
        return []
    
    cleaned_table = []
    for row in table:
        if row:  # Skip empty rows
            cleaned_row = []
            for cell in row:
                if cell:
                    # Clean cell content
                    cell_str = str(cell).strip()
                    # Remove excessive whitespace
                    cell_str = ' '.join(cell_str.split())
                    cleaned_row.append(cell_str)
                else:
                    cleaned_row.append('')
            
            # Only add row if it has meaningful content
            if any(cell for cell in cleaned_row):
                cleaned_table.append(cleaned_row)
    
    return cleaned_table

def is_duplicate_table(new_table, existing_tables, similarity_threshold=0.8):
    """
    Check if a table is a duplicate of existing tables
    """
    if not existing_tables:
        return False
    
    for existing_table in existing_tables:
        if len(new_table) == len(existing_table):
            matching_rows = 0
            for i, row in enumerate(new_table):
                if i < len(existing_table) and row == existing_table[i]:
                    matching_rows += 1
            
            similarity = matching_rows / len(new_table) if new_table else 0
            if similarity >= similarity_threshold:
                return True
    
    return False

def add_table_borders(worksheet, rows, cols):
    """
    Add borders to table cells
    """
    from openpyxl.styles import Border, Side
    
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    for row in range(1, rows + 1):
        for col in range(1, cols + 1):
            worksheet.cell(row=row, column=col).border = thin_border

def adjust_column_widths(worksheet, table_data):
    """
    Smart column width adjustment
    """
    from openpyxl.utils import get_column_letter
    
    if not table_data:
        return
    
    for col in range(1, len(table_data[0]) + 1):
        max_length = 0
        for row in range(1, len(table_data) + 1):
            cell_value = worksheet.cell(row=row, column=col).value
            if cell_value:
                # Count characters, considering some characters are wider
                cell_length = len(str(cell_value))
                # Adjust for special characters
                for char in str(cell_value):
                    if char in 'WMmw':
                        cell_length += 0.5
                    elif char in 'il':
                        cell_length -= 0.2
                
                max_length = max(max_length, cell_length)
        
        # Set column width with reasonable limits
        adjusted_width = min(max(max_length + 2, 8), 50)
        worksheet.column_dimensions[get_column_letter(col)].width = adjusted_width

def parse_page_selection(page_selection):
    """
    Parse page selection string into list of page numbers
    Examples: "1,3-5,7" -> [1,3,4,5,7]
             "all" or "" -> None (all pages)
    """
    if not page_selection or page_selection.lower() == 'all':
        return None  # All pages
    
    pages = set()
    parts = page_selection.split(',')
    
    for part in parts:
        part = part.strip()
        if '-' in part:
            # Handle range like "3-5"
            try:
                start, end = map(int, part.split('-'))
                pages.update(range(start, end + 1))
            except ValueError:
                print(f"Invalid range format: {part}")
        else:
            # Handle single page
            try:
                pages.add(int(part))
            except ValueError:
                print(f"Invalid page number: {part}")
    
    return sorted(list(pages)) if pages else None

def pdf_to_powerpoint(pdf_file_path, pptx_file_path):
    """
    Convert PDF to PowerPoint using pdf2image and python-pptx
    Converts each PDF page to an image and creates a PowerPoint slide for each page
    """
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import tempfile
        import os
        
        print(f"Starting conversion of {pdf_file_path} to {pptx_file_path}")
        
        # Convert PDF pages to images with high quality
        print("Converting PDF pages to images...")
        images = convert_from_path(
            pdf_file_path, 
            dpi=200,  # High resolution for better quality
            fmt='PNG',  # PNG format for better quality
            thread_count=4  # Use multiple threads for faster processing
        )
        
        print(f"Converted {len(images)} pages to images")
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Process each page
        for page_num, image in enumerate(images):
            print(f"Processing page {page_num + 1}/{len(images)}")
            
            # Add slide with blank layout
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Save image to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                image.save(temp_img.name, 'PNG')
                temp_img_path = temp_img.name
            
            try:
                # Add image to slide
                # Calculate dimensions to fit slide while maintaining aspect ratio
                slide_width = Inches(10)
                slide_height = Inches(7.5)
                
                # Get image dimensions
                img_width, img_height = image.size
                
                # Calculate scaling to fit slide
                width_ratio = slide_width.inches / (img_width / 200)  # 200 DPI
                height_ratio = slide_height.inches / (img_height / 200)
                scale_ratio = min(width_ratio, height_ratio)
                
                # Calculate final dimensions
                final_width = Inches((img_width / 200) * scale_ratio)
                final_height = Inches((img_height / 200) * scale_ratio)
                
                # Center the image on slide
                left = (slide_width - final_width) / 2
                top = (slide_height - final_height) / 2
                
                # Add image to slide
                slide.shapes.add_picture(temp_img_path, left, top, final_width, final_height)
                
                print(f"Added page {page_num + 1} to slide {page_num + 1}")
                
            finally:
                # Clean up temporary image file
                if os.path.exists(temp_img_path):
                    os.unlink(temp_img_path)
        
        # Save the presentation
        prs.save(pptx_file_path)
        print(f"Successfully created PowerPoint file: {pptx_file_path}")
        return True
        
    except Exception as e:
        print(f"Error converting PDF to PowerPoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def pdf_to_powerpoint_text(pdf_file_path, pptx_file_path):
    """
    Alternative method: Convert PDF to PowerPoint using text extraction
    Creates text-based slides from PDF content
    """
    try:
        import pdfplumber
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        
        print(f"Starting text-based conversion of {pdf_file_path} to {pptx_file_path}")
        
        # Extract text from PDF
        with pdfplumber.open(pdf_file_path) as pdf:
            prs = Presentation()
            
            for page_num, page in enumerate(pdf.pages):
                print(f"Processing page {page_num + 1}/{len(pdf.pages)}")
                
                # Extract text
                text = page.extract_text()
                
                if text and text.strip():
                    # Add slide with title and content layout
                    slide_layout = prs.slide_layouts[1]  # Title and content
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = f"Page {page_num + 1}"
                    
                    # Add content
                    content = slide.placeholders[1]
                    content.text = text.strip()
                    
                    print(f"Added text from page {page_num + 1} to slide {page_num + 1}")
                else:
                    # If no text, create blank slide with page number
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add text box with page number
                    left = Inches(1)
                    top = Inches(1)
                    width = Inches(8)
                    height = Inches(1)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = f"Page {page_num + 1} (No text content)"
                    
                    print(f"Added blank slide for page {page_num + 1}")
            
            # Save the presentation
            prs.save(pptx_file_path)
            print(f"Successfully created text-based PowerPoint file: {pptx_file_path}")
            return True
            
    except Exception as e:
        print(f"Error in text-based PDF to PowerPoint conversion: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def pdf_to_text(input_path, output_path, options=None):
    """
    Convert PDF to text file using pdfplumber with OCR support
    """
    try:
        import pdfplumber
        import re
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image
        import tempfile
        import os
        
        if options is None:
            options = {}
        
        use_ocr = options.get('ocr', True)
        layout_mode = options.get('layout', 'formatted')  # Changed default to formatted
        
        extracted_text = []
        
        with pdfplumber.open(input_path) as pdf:
            total_pages = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                print(f"Processing page {page_num}/{total_pages}")
                
                # Extract text from the page
                text = page.extract_text()
                
                if not text and use_ocr:
                    # If no text found and OCR is enabled, try OCR
                    print(f"Page {page_num}: No text found, attempting OCR...")
                    try:
                        # Convert PDF page to image for OCR
                        images = convert_from_path(input_path, first_page=page_num, last_page=page_num)
                        if images:
                            # Get the first (and only) image
                            image = images[0]
                            
                            # Perform OCR on the image
                            ocr_text = pytesseract.image_to_string(image, lang='eng')
                            
                            if ocr_text.strip():
                                text = ocr_text
                                print(f"Page {page_num}: OCR successful, extracted {len(text)} characters")
                            else:
                                text = f"[OCR completed but no text found on page {page_num}]"
                                print(f"Page {page_num}: OCR completed but no text found")
                        else:
                            text = f"[OCR failed - could not convert page {page_num} to image]"
                            print(f"Page {page_num}: OCR failed - could not convert to image")
                    except Exception as ocr_error:
                        print(f"Page {page_num}: OCR error: {str(ocr_error)}")
                        text = f"[OCR error on page {page_num}: {str(ocr_error)}]"
                
                if text:
                    if layout_mode == 'formatted':
                        # Try to preserve some formatting
                        text = re.sub(r'\n{3,}', '\n\n', text)  # Remove excessive line breaks
                        text = re.sub(r' +', ' ', text)  # Normalize spaces
                        # Keep paragraph breaks and basic formatting
                        text = text.strip()
                    else:
                        # Simple text flow - remove extra whitespace
                        text = re.sub(r'\s+', ' ', text).strip()
                    
                    extracted_text.append(text)
                else:
                    extracted_text.append(f"[Page {page_num}: No text content]")
        
        # Combine all text
        full_text = '\n\n'.join(extracted_text)
        
        # Write to output file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        
        print(f"Text extraction completed. Output saved to: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error in pdf_to_text: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def pdf_to_html(input_path, output_path, options=None):
    """
    Convert PDF to HTML file using pdfplumber with image embedding support
    """
    try:
        import pdfplumber
        import re
        import base64
        import io
        from PIL import Image
        from pdf2image import convert_from_path
        import pytesseract
        
        if options is None:
            options = {}
        
        use_ocr = options.get('ocr', False)
        embed_images = options.get('embedImages', True)
        responsive = options.get('responsive', False)
        
        html_parts = []
        
        # HTML header with responsive design if enabled
        if responsive:
            html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Converted PDF</title>
    <style>
        body { font-family: Arial, sans-serif; line-height: 1.6; margin: 20px; }
        .page { margin-bottom: 30px; page-break-after: always; }
        .page:last-child { page-break-after: auto; }
        .page-number { font-weight: bold; color: #666; margin-bottom: 10px; }
        img { max-width: 100%; height: auto; }
        @media print { .page { page-break-after: always; } }
        @media screen and (max-width: 768px) { body { margin: 10px; } }
    </style>
</head>
<body>""")
        else:
            html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Converted PDF</title>
    <style>
        body { font-family: Arial, sans-serif; line-height: 1.6; margin: 20px; }
        .page { margin-bottom: 30px; }
        .page-number { font-weight: bold; color: #666; margin-bottom: 10px; }
        img { max-width: 800px; }
    </style>
</head>
<body>""")
        
        with pdfplumber.open(input_path) as pdf:
            total_pages = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                print(f"Processing page {page_num}/{total_pages}")
                
                # Start page div
                html_parts.append(f'<div class="page">')
                html_parts.append(f'<div class="page-number">Page {page_num}</div>')
                
                # Extract text from the page
                text = page.extract_text()
                
                if not text and use_ocr:
                    # If no text found and OCR is enabled, try OCR
                    print(f"Page {page_num}: No text found, attempting OCR...")
                    try:
                        # Convert PDF page to image for OCR
                        images = convert_from_path(input_path, first_page=page_num, last_page=page_num)
                        if images:
                            image = images[0]
                            ocr_text = pytesseract.image_to_string(image, lang='eng')
                            if ocr_text.strip():
                                text = ocr_text
                                print(f"Page {page_num}: OCR successful")
                            else:
                                text = f"[OCR completed but no text found on page {page_num}]"
                        else:
                            text = f"[OCR failed - could not convert page {page_num} to image]"
                    except Exception as ocr_error:
                        print(f"Page {page_num}: OCR error: {str(ocr_error)}")
                        text = f"[OCR error on page {page_num}]"
                
                # Add text content
                if text:
                    # Convert text to HTML paragraphs
                    paragraphs = text.split('\n\n')
                    for para in paragraphs:
                        if para.strip():
                            # Escape HTML characters
                            para_html = para.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            html_parts.append(f'<p>{para_html}</p>')
                else:
                    html_parts.append(f'<p>[No text content on page {page_num}]</p>')
                
                # Extract and embed images if enabled
                if embed_images:
                    try:
                        # Convert page to image
                        images = convert_from_path(input_path, first_page=page_num, last_page=page_num)
                        if images:
                            image = images[0]
                            
                            # Convert PIL image to base64
                            img_buffer = io.BytesIO()
                            image.save(img_buffer, format='PNG')
                            img_str = base64.b64encode(img_buffer.getvalue()).decode()
                            
                            # Add image to HTML
                            html_parts.append(f'<img src="data:image/png;base64,{img_str}" alt="Page {page_num}" />')
                            print(f"Page {page_num}: Image embedded")
                    except Exception as img_error:
                        print(f"Page {page_num}: Image embedding error: {str(img_error)}")
                        html_parts.append(f'<p>[Image embedding failed for page {page_num}]</p>')
                
                # End page div
                html_parts.append('</div>')
        
        # HTML footer
        html_parts.append("</body></html>")
        
        # Combine all HTML parts
        full_html = '\n'.join(html_parts)
        
        # Write to output file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_html)
        
        print(f"HTML conversion completed. Output saved to: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error in pdf_to_html: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def pdf_to_epub(pdf_file_path, epub_file_path, options=None):
    """
    Convert PDF to EPUB format using pypdf2 and ebooklib with advanced options
    """
    try:
        import PyPDF2
        from ebooklib import epub
        import re
        import html
        import base64
        import io
        from PIL import Image
        
        print(f"Starting conversion of {pdf_file_path} to {epub_file_path}")
        
        if options is None:
            options = {}
        
        # Create EPUB book
        book = epub.EpubBook()
        
        # Set metadata
        custom_title = options.get('custom_title')
        custom_author = options.get('custom_author')
        
        book.set_identifier(f"pdf-conversion-{os.path.basename(pdf_file_path)}")
        book.set_title(custom_title or f"Converted from {os.path.basename(pdf_file_path)}")
        book.set_language('en')
        book.add_author(custom_author or 'PDF Converter')
        
        # Extract text and images from PDF
        chapters = []
        chapter_titles = []
        images = []
        
        with open(pdf_file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            print(f"Processing {total_pages} pages")
            
            # Determine chapter size (pages per chapter)
            pages_per_chapter = options.get('pages_per_chapter', 10)
            if pages_per_chapter <= 0:
                pages_per_chapter = total_pages  # Single chapter
            
            current_chapter_text = []
            current_chapter_pages = []
            page_break_style = options.get('page_break_style', 'chapter')
            include_images = options.get('include_images', True)
            image_quality = options.get('image_quality', 'medium')
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                print(f"Processing page {page_num}/{total_pages}")
                
                # Extract text from page
                text = page.extract_text()
                
                # If no text found and OCR is enabled, try OCR
                if not text and options.get('ocr', False):
                    print(f"Page {page_num}: No text found, attempting OCR...")
                    try:
                        from pdf2image import convert_from_path
                        import pytesseract
                        
                        # Convert PDF page to image for OCR
                        images_pdf = convert_from_path(pdf_file_path, first_page=page_num, last_page=page_num)
                        if images_pdf:
                            image = images_pdf[0]
                            ocr_text = pytesseract.image_to_string(image, lang='eng')
                            if ocr_text.strip():
                                text = ocr_text
                                print(f"Page {page_num}: OCR successful")
                            else:
                                text = f"[OCR completed but no text found on page {page_num}]"
                        else:
                            text = f"[OCR failed - could not convert page {page_num} to image]"
                    except Exception as ocr_error:
                        print(f"Page {page_num}: OCR error: {str(ocr_error)}")
                        text = f"[OCR error on page {page_num}]"
                
                # Extract images if enabled
                page_images = []
                if include_images:
                    try:
                        from pdf2image import convert_from_path
                        images_pdf = convert_from_path(pdf_file_path, first_page=page_num, last_page=page_num)
                        if images_pdf:
                            image = images_pdf[0]
                            
                            # Adjust image quality based on setting
                            if image_quality == 'low':
                                image = image.resize((image.width // 2, image.height // 2), Image.Resampling.LANCZOS)
                            elif image_quality == 'high':
                                # Keep original size for high quality
                                pass
                            else:  # medium
                                # Slight reduction for medium quality
                                image = image.resize((int(image.width * 0.8), int(image.height * 0.8)), Image.Resampling.LANCZOS)
                            
                            # Convert to base64
                            img_buffer = io.BytesIO()
                            image.save(img_buffer, format='JPEG', quality=85)
                            img_str = base64.b64encode(img_buffer.getvalue()).decode()
                            
                            page_images.append({
                                'data': img_str,
                                'alt': f'Page {page_num}'
                            })
                            print(f"Page {page_num}: Image extracted")
                    except Exception as img_error:
                        print(f"Page {page_num}: Image extraction error: {str(img_error)}")
                
                if text:
                    # Clean and format text
                    cleaned_text = clean_text_for_epub(text, options)
                    current_chapter_text.append(cleaned_text)
                    current_chapter_pages.append(page_num)
                    
                    # Add images to chapter if any
                    if page_images:
                        for img in page_images:
                            current_chapter_text.append(f'<img src="data:image/jpeg;base64,{img["data"]}" alt="{img["alt"]}" style="max-width: 100%; height: auto;" />')
                
                # Create new chapter based on page break style
                should_create_chapter = False
                if page_break_style == 'page':
                    should_create_chapter = True
                elif page_break_style == 'chapter':
                    should_create_chapter = (len(current_chapter_pages) >= pages_per_chapter or page_num == total_pages)
                else:  # none
                    should_create_chapter = (page_num == total_pages)
                
                if should_create_chapter and current_chapter_text:
                    # Create chapter content
                    chapter_content = '\n\n'.join(current_chapter_text)
                    
                    # Create chapter title
                    if len(current_chapter_pages) == 1:
                        chapter_title = f"Page {current_chapter_pages[0]}"
                    else:
                        chapter_title = f"Pages {current_chapter_pages[0]}-{current_chapter_pages[-1]}"
                    
                    # Create EPUB chapter
                    chapter = epub.EpubHtml(
                        title=chapter_title,
                        file_name=f'chapter_{len(chapters) + 1}.xhtml',
                        content=f'<h1>{chapter_title}</h1>\n{chapter_content}'
                    )
                    
                    chapters.append(chapter)
                    chapter_titles.append(chapter_title)
                    
                    # Reset for next chapter
                    current_chapter_text = []
                    current_chapter_pages = []
            
            # Add chapters to book
            for chapter in chapters:
                book.add_item(chapter)
            
            # Create table of contents if enabled
            if options.get('add_toc', True):
                book.toc = [(epub.Section('Chapters'), chapters)]
            
            # Add default NCX and Nav files
            book.add_item(epub.EpubNcx())
            book.add_item(epub.EpubNav())
            
            # Define CSS style based on options
            font_size_map = {'small': '12px', 'medium': '16px', 'large': '20px'}
            line_spacing_map = {'tight': '1.2', 'normal': '1.5', 'loose': '1.8'}
            
            font_size = font_size_map.get(options.get('font_size', 'medium'), '16px')
            line_spacing = line_spacing_map.get(options.get('line_spacing', 'normal'), '1.5')
            
            style = f'''
            @namespace epub "http://www.idpf.org/2007/ops";
            body {{ 
                font-family: Arial, sans-serif; 
                line-height: {line_spacing}; 
                margin: 20px; 
                font-size: {font_size};
            }}
            h1 {{ 
                color: #333; 
                border-bottom: 2px solid #333; 
                padding-bottom: 10px; 
                font-size: 1.5em;
            }}
            p {{ 
                margin-bottom: 1em; 
                text-align: justify; 
            }}
            img {{ 
                max-width: 100%; 
                height: auto; 
                display: block; 
                margin: 1em auto; 
            }}
            '''
            
            # Add CSS file
            nav_css = epub.EpubItem(
                uid="style_nav",
                file_name="style/nav.css",
                media_type="text/css",
                content=style
            )
            book.add_item(nav_css)
            
            # Set CSS for chapters
            for chapter in chapters:
                chapter.add_item(nav_css)
            
            # Create spine
            book.spine = ['nav'] + chapters
            
            # Write EPUB file
            epub.write_epub(epub_file_path, book)
            
            print(f"EPUB conversion completed. Output saved to: {epub_file_path}")
            return True
            
    except Exception as e:
        print(f"Error in pdf_to_epub: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def clean_text_for_epub(text, options=None):
    """
    Clean and format text for EPUB conversion with advanced options
    """
    import re
    import html
    
    if not text:
        return ""
    
    preserve_layout = options.get('preserve_layout', False) if options else False
    
    if preserve_layout:
        # Preserve more of the original layout
        # Split by newlines but keep more structure
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if line:
                # Escape HTML characters
                line = html.escape(line)
                # Keep line breaks for layout preservation
                cleaned_lines.append(f'<p>{line}</p>')
        return '\n'.join(cleaned_lines)
    else:
        # Standard text cleaning
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Split into paragraphs
        paragraphs = text.split('\n')
        
        # Clean each paragraph
        cleaned_paragraphs = []
        for para in paragraphs:
            para = para.strip()
            if para:
                # Escape HTML characters
                para = html.escape(para)
                # Convert to HTML paragraph
                cleaned_paragraphs.append(f'<p>{para}</p>')
        
        return '\n'.join(cleaned_paragraphs)

def pdf_to_rtf(pdf_file_path, rtf_file_path, options=None):
    """
    Convert PDF to RTF format using pypdf2 and custom RTF generation
    """
    try:
        import PyPDF2
        import re
        import base64
        import io
        from PIL import Image
        
        print(f"Starting conversion of {pdf_file_path} to {rtf_file_path}")
        
        if options is None:
            options = {}
        
        # Parse page selection
        page_selection = options.get('page_selection', 'all')
        selected_pages = parse_page_selection(page_selection)
        
        # Get options
        preserve_formatting = options.get('preserve_formatting', True)
        include_images = options.get('include_images', False)
        font_size_map = {'small': '10', 'medium': '12', 'large': '14'}
        line_spacing_map = {'single': '1', 'normal': '1.15', 'double': '2'}
        
        font_size = font_size_map.get(options.get('font_size', 'medium'), '12')
        line_spacing = line_spacing_map.get(options.get('line_spacing', 'normal'), '1.15')
        add_page_breaks = options.get('page_breaks', True)
        custom_title = options.get('custom_title', '')
        
        # RTF header
        rtf_content = []
        rtf_content.append(r'{\rtf1\ansi\deff0')
        rtf_content.append(r'{\fonttbl{\f0\froman\fcharset0 Times New Roman;}}')
        rtf_content.append(r'{\colortbl ;\red0\green0\blue0;}')
        rtf_content.append(r'{\*\generator PDF to RTF Converter;}')
        
        # Document properties
        if custom_title:
            rtf_content.append(f'{{\\*\\title {escape_rtf_text(custom_title)}}}')
        
        # Default formatting
        rtf_content.append(f'\\f0\\fs{int(font_size) * 2}')
        rtf_content.append(f'\\sl{int(float(line_spacing) * 240)}')
        
        with open(pdf_file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            print(f"Processing {total_pages} pages")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                # Skip pages not in selection
                if selected_pages and page_num not in selected_pages:
                    continue
                
                print(f"Processing page {page_num}/{total_pages}")
                
                # Extract text from page
                text = page.extract_text()
                
                # If no text found and OCR is enabled, try OCR
                if not text and options.get('ocr', False):
                    print(f"Page {page_num}: No text found, attempting OCR...")
                    try:
                        from pdf2image import convert_from_path
                        import pytesseract
                        
                        # Convert PDF page to image for OCR
                        images_pdf = convert_from_path(pdf_file_path, first_page=page_num, last_page=page_num)
                        if images_pdf:
                            image = images_pdf[0]
                            ocr_text = pytesseract.image_to_string(image, lang='eng')
                            if ocr_text.strip():
                                text = ocr_text
                                print(f"Page {page_num}: OCR successful")
                            else:
                                text = f"[OCR completed but no text found on page {page_num}]"
                        else:
                            text = f"[OCR failed - could not convert page {page_num} to image]"
                    except Exception as ocr_error:
                        print(f"Page {page_num}: OCR error: {str(ocr_error)}")
                        text = f"[OCR error on page {page_num}]"
                
                # Process text
                if text:
                    # Clean and format text
                    if preserve_formatting:
                        # Preserve more formatting
                        formatted_text = format_text_for_rtf_preserved(text)
                    else:
                        # Simple formatting
                        formatted_text = format_text_for_rtf_simple(text)
                    
                    rtf_content.append(formatted_text)
                
                # Extract images if enabled
                if include_images:
                    try:
                        from pdf2image import convert_from_path
                        images_pdf = convert_from_path(pdf_file_path, first_page=page_num, last_page=page_num)
                        if images_pdf:
                            image = images_pdf[0]
                            
                            # Resize image for RTF (RTF has size limitations)
                            max_width = 400
                            if image.width > max_width:
                                ratio = max_width / image.width
                                new_height = int(image.height * ratio)
                                image = image.resize((max_width, new_height), Image.Resampling.LANCZOS)
                            
                            # Convert to JPEG for better RTF compatibility
                            img_buffer = io.BytesIO()
                            image.save(img_buffer, format='JPEG', quality=85)
                            img_data = img_buffer.getvalue()
                            
                            # Convert to hex for RTF (RTF uses hex encoding, not base64)
                            hex_data = ''.join([f'{b:02x}' for b in img_data])
                            
                            # Add image to RTF with proper format
                            rtf_content.append('\\par')  # Add space before image
                            rtf_content.append(f'{{\\*\\shppict{{\\pict\\jpegblip\\picwgoal{image.width * 15}\\pichgoal{image.height * 15}')
                            rtf_content.append(hex_data)
                            rtf_content.append('}}}')
                            rtf_content.append('\\par')  # Add space after image
                            print(f"Page {page_num}: Image added to RTF (size: {image.width}x{image.height})")
                    except Exception as img_error:
                        print(f"Page {page_num}: Image extraction error: {str(img_error)}")
                        import traceback
                        traceback.print_exc()
                
                # Add page break if enabled and not the last page
                if add_page_breaks and page_num < total_pages:
                    rtf_content.append('\\page')
        
        # RTF footer
        rtf_content.append('}')
        
        # Write RTF file
        with open(rtf_file_path, 'w', encoding='utf-8') as f:
            f.write(''.join(rtf_content))
        
        print(f"RTF conversion completed. Output saved to: {rtf_file_path}")
        return True
        
    except Exception as e:
        print(f"Error in pdf_to_rtf: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def escape_rtf_text(text):
    """
    Escape text for RTF format
    """
    if not text:
        return ""
    
    # RTF escape sequences
    text = text.replace('\\', '\\\\')
    text = text.replace('{', '\\{')
    text = text.replace('}', '\\}')
    text = text.replace('\n', '\\par ')
    text = text.replace('\r', '\\par ')
    text = text.replace('\t', '\\tab ')
    
    return text

def format_text_for_rtf_simple(text):
    """
    Format text for RTF with simple formatting
    """
    if not text:
        return ""
    
    # Clean text
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # Escape for RTF
    text = escape_rtf_text(text)
    
    # Add paragraph formatting
    return f'\\par {text}'

def format_text_for_rtf_preserved(text):
    """
    Format text for RTF with preserved formatting
    """
    if not text:
        return ""
    
    # Split into lines to preserve structure
    lines = text.split('\n')
    formatted_lines = []
    
    for line in lines:
        line = line.strip()
        if line:
            # Escape for RTF
            escaped_line = escape_rtf_text(line)
            formatted_lines.append(f'\\par {escaped_line}')
        else:
            # Empty line
            formatted_lines.append('\\par ')
    
    return ''.join(formatted_lines)

def pdf_to_svg(pdf_file_path, svg_file_path, options=None):
    """
    Convert PDF to SVG using pdf2image
    """
    try:
        from pdf2image import convert_from_path
        import io
        from PIL import Image
        import os
        
        print(f"Starting conversion of {pdf_file_path} to {svg_file_path}")
        
        # Get options with defaults
        if options is None:
            options = {}
        
        dpi = options.get('dpi', 300)
        width = options.get('width', 800)
        height = options.get('height', 600)
        page_selection = options.get('page_selection', 'all')
        output_mode = options.get('output_mode', 'single')  # 'single', 'per_page', 'grouped'
        pages_per_svg = options.get('pages_per_svg', 1)  # For grouped mode
        
        print(f"Conversion options: DPI={dpi}, Width={width}, Height={height}, Pages={page_selection}, Mode={output_mode}, PagesPerSVG={pages_per_svg}")
        
        # Convert PDF pages to images first to get total pages
        pages = convert_from_path(pdf_file_path, dpi=dpi)
        total_pages = len(pages)
        print(f"Total pages in PDF: {total_pages}")
        
        # Validate page selection for SVG conversion only
        selected_pages, error_message = validate_page_selection_for_svg(page_selection, total_pages)
        if error_message:
            print(f"Page selection error: {error_message}")
            return False
        
        # Filter pages if selection specified
        if selected_pages:
            pages = [pages[i-1] for i in selected_pages if 1 <= i <= len(pages)]
            print(f"Selected pages: {selected_pages}")
        
        if not pages:
            print("No pages to convert")
            return False
        
        # Get base filename without extension
        base_filename = os.path.splitext(os.path.basename(svg_file_path))[0]
        output_dir = os.path.dirname(svg_file_path)
        
        if output_mode == 'per_page':
            # Create separate SVG file for each page
            print(f"Creating {len(pages)} separate SVG files")
            for i, page in enumerate(pages):
                page_filename = f"{base_filename}_page_{i+1}.svg"
                page_path = os.path.join(output_dir, page_filename)
                
                # Process single page
                svg_content = create_svg_content([page], width, height, [f"Page {i+1}"])
                
                # Write SVG file
                with open(page_path, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(svg_content))
                
                print(f"Created SVG file: {page_filename}")
            
            # Create a summary file that lists all created files
            summary_path = os.path.join(output_dir, f"{base_filename}_summary.txt")
            with open(summary_path, 'w', encoding='utf-8') as f:
                f.write(f"PDF to SVG conversion summary\n")
                f.write(f"Original PDF: {os.path.basename(pdf_file_path)}\n")
                f.write(f"Total pages converted: {len(pages)}\n")
                f.write(f"Output mode: Per-page SVG files\n\n")
                f.write("Created files:\n")
                for i in range(len(pages)):
                    f.write(f"- {base_filename}_page_{i+1}.svg\n")
            
            print(f"Created summary file: {os.path.basename(summary_path)}")
            return True
            
        elif output_mode == 'grouped':
            # Group pages into multiple SVG files
            total_groups = (len(pages) + pages_per_svg - 1) // pages_per_svg
            print(f"Creating {total_groups} SVG files with {pages_per_svg} pages each")
            
            for group_idx in range(total_groups):
                start_idx = group_idx * pages_per_svg
                end_idx = min(start_idx + pages_per_svg, len(pages))
                group_pages = pages[start_idx:end_idx]
                
                group_filename = f"{base_filename}_group_{group_idx + 1}.svg"
                group_path = os.path.join(output_dir, group_filename)
                
                # Create page labels for this group
                page_labels = [f"Page {start_idx + i + 1}" for i in range(len(group_pages))]
                
                # Process group of pages
                svg_content = create_svg_content(group_pages, width, height, page_labels)
                
                # Write SVG file
                with open(group_path, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(svg_content))
                
                print(f"Created SVG file: {group_filename} (pages {start_idx + 1}-{end_idx})")
            
            # Create a summary file
            summary_path = os.path.join(output_dir, f"{base_filename}_summary.txt")
            with open(summary_path, 'w', encoding='utf-8') as f:
                f.write(f"PDF to SVG conversion summary\n")
                f.write(f"Original PDF: {os.path.basename(pdf_file_path)}\n")
                f.write(f"Total pages: {len(pages)}\n")
                f.write(f"Output mode: Grouped ({pages_per_svg} pages per SVG)\n")
                f.write(f"Total SVG files created: {total_groups}\n\n")
                f.write("Created files:\n")
                for group_idx in range(total_groups):
                    start_idx = group_idx * pages_per_svg
                    end_idx = min(start_idx + pages_per_svg, len(pages))
                    f.write(f"- {base_filename}_group_{group_idx + 1}.svg (pages {start_idx + 1}-{end_idx})\n")
            
            print(f"Created summary file: {os.path.basename(summary_path)}")
            return True
            
        else:
            # Default: Single SVG file with all pages
            print("Creating single SVG file with all pages")
            page_labels = [f"Page {i+1}" for i in range(len(pages))]
            svg_content = create_svg_content(pages, width, height, page_labels)
            
            # Write SVG file
            with open(svg_file_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(svg_content))
            
            print(f"Successfully converted {pdf_file_path} to {svg_file_path}")
            print(f"SVG contains {len(pages)} pages")
            return True
        
    except Exception as e:
        print(f"Error converting PDF to SVG: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def validate_page_selection_for_svg(page_selection, total_pages):
    """
    Validate page selection for SVG conversion only
    Returns: (pages_list, error_message) where error_message is None if valid
    """
    if not page_selection or page_selection.lower() == 'all':
        return None, None  # All pages
    
    pages = set()
    parts = page_selection.split(',')
    
    for part in parts:
        part = part.strip()
        if '-' in part:
            # Handle range like "3-5"
            try:
                start, end = map(int, part.split('-'))
                if start < 1:
                    return None, f"Invalid page range '{part}': page numbers must be 1 or greater"
                if end > total_pages:
                    return None, f"Invalid page range '{part}': page {end} exceeds total pages ({total_pages})"
                if start > end:
                    return None, f"Invalid page range '{part}': start page ({start}) is greater than end page ({end})"
                pages.update(range(start, end + 1))
            except ValueError:
                return None, f"Invalid range format: '{part}'. Use format like '1-3'"
        else:
            # Handle single page
            try:
                page_num = int(part)
                if page_num < 1:
                    return None, f"Invalid page number '{part}': page numbers must be 1 or greater"
                if page_num > total_pages:
                    return None, f"Invalid page number '{part}': page {page_num} exceeds total pages ({total_pages})"
                pages.add(page_num)
            except ValueError:
                return None, f"Invalid page number: '{part}'. Use numbers only"
    
    return sorted(list(pages)) if pages else None, None

def create_svg_content(pages, width, height, page_labels):
    """
    Create SVG content from a list of pages
    """
    from PIL import Image
    import io
    
    svg_content = []
    svg_content.append('<?xml version="1.0" encoding="UTF-8"?>')
    svg_content.append('<svg xmlns="http://www.w3.org/2000/svg" version="1.1" width="100%" height="100%">')
    
    # Add metadata
    svg_content.append('<defs>')
    svg_content.append('<style>')
    svg_content.append('  .page { margin: 10px; }')
    svg_content.append('  .page-image { max-width: 100%; height: auto; }')
    svg_content.append('</style>')
    svg_content.append('</defs>')
    
    current_y = 0
    
    for i, page in enumerate(pages):
        print(f"Processing {page_labels[i]}")
        
        # Resize image if needed
        if page.width > width or page.height > height:
            ratio = min(width / page.width, height / page.height)
            new_width = int(page.width * ratio)
            new_height = int(page.height * ratio)
            page = page.resize((new_width, new_height), Image.Resampling.LANCZOS)
            print(f"Resized {page_labels[i]} to {new_width}x{new_height}")
        
        # Convert image to base64
        img_buffer = io.BytesIO()
        page.save(img_buffer, format='PNG')
        img_data = img_buffer.getvalue()
        import base64
        base64_data = base64.b64encode(img_data).decode('utf-8')
        
        # Add page as SVG element
        svg_content.append(f'<g id="{page_labels[i].replace(" ", "_").lower()}" class="page" transform="translate(0, {current_y})">')
        svg_content.append(f'<image x="0" y="0" width="{page.width}" height="{page.height}" class="page-image"')
        svg_content.append(f'href="data:image/png;base64,{base64_data}"/>')
        svg_content.append('</g>')
        
        current_y += page.height + 20  # Add spacing between pages
    
    svg_content.append('</svg>')
    return svg_content

def merge_pdfs(output_path, input_paths):
    """
    Merges multiple PDF files into a single PDF.
    """
    try:
        import signal
        
        # Set up timeout handler
        def timeout_handler(signum, frame):
            raise TimeoutError("PDF merge operation timed out")
        
        # Set timeout for 4 minutes
        signal.signal(signal.SIGALRM, timeout_handler)
        signal.alarm(240)  # 4 minutes
        
        try:
            from PyPDF2 import PdfMerger
            
            merger = PdfMerger()
            
            print(f"Starting merge process for {len(input_paths)} files.")
            for pdf_path in input_paths:
                if os.path.exists(pdf_path):
                    print(f"Appending file: {pdf_path}")
                    merger.append(pdf_path)
                else:
                    print(f"Warning: File not found and skipped: {pdf_path}")
            
            # Write out the merged PDF
            merger.write(output_path)
            merger.close()
            
            print(f"Successfully merged {len(input_paths)} files into {output_path}")
            return True
        finally:
            # Cancel the alarm
            signal.alarm(0)
            
    except TimeoutError as e:
        print(f"PDF merge operation timed out: {str(e)}")
        return False
    except Exception as e:
        print(f"Error in merge_pdfs: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def split_pdf(input_path, output_zip_path, options=None):
    """
    Splits a PDF into multiple files based on specified ranges or settings.
    The output is a ZIP file containing the split PDFs.
    """
    try:
        import signal
        
        # Set up timeout handler
        def timeout_handler(signum, frame):
            raise TimeoutError("PDF split operation timed out")
        
        # Set timeout for 4 minutes
        signal.signal(signal.SIGALRM, timeout_handler)
        signal.alarm(240)  # 4 minutes
        
        try:
            import PyPDF2
            import zipfile
            import tempfile
            import re

            if options is None:
                options = {}
            
            split_mode = options.get('split_mode', 'ranges')
            page_ranges_str = options.get('page_ranges', '')

            pdf_reader = PyPDF2.PdfReader(input_path)
            total_pages = len(pdf_reader.pages)
            
            # Create a temporary directory to store the split PDF files before zipping
            with tempfile.TemporaryDirectory() as temp_dir:
                split_file_paths = []

                if split_mode == 'all':
                    # Split every page into a separate file
                    for i in range(total_pages):
                        pdf_writer = PyPDF2.PdfWriter()
                        pdf_writer.add_page(pdf_reader.pages[i])
                        output_filename = os.path.join(temp_dir, f'page_{i + 1}.pdf')
                        with open(output_filename, 'wb') as out_f:
                            pdf_writer.write(out_f)
                        split_file_paths.append(output_filename)
                    print(f"Split PDF into {total_pages} separate pages.")

                elif split_mode == 'ranges':
                    # Split by custom page ranges
                    if not page_ranges_str:
                        raise ValueError("Page ranges are required for 'ranges' split mode.")
                    
                    range_groups = page_ranges_str.split(',')
                    for i, group in enumerate(range_groups):
                        group = group.strip()
                        pdf_writer = PyPDF2.PdfWriter()
                        
                        # Parse individual pages and ranges (e.g., "1-3" or "5")
                        pages_in_group = set()
                        parts = re.split(r'[-–]', group) # Handles hyphen and en-dash
                        if len(parts) == 1:
                            page_num = int(parts[0])
                            if 1 <= page_num <= total_pages:
                                pages_in_group.add(page_num)
                        elif len(parts) == 2:
                            start, end = int(parts[0]), int(parts[1])
                            for page_num in range(start, end + 1):
                                if 1 <= page_num <= total_pages:
                                    pages_in_group.add(page_num)

                        if not pages_in_group:
                            continue # Skip empty or invalid groups

                        # Add sorted pages to the writer
                        for page_num in sorted(list(pages_in_group)):
                            pdf_writer.add_page(pdf_reader.pages[page_num - 1])
                        
                        output_filename = os.path.join(temp_dir, f'split_group_{i + 1}.pdf')
                        with open(output_filename, 'wb') as out_f:
                            pdf_writer.write(out_f)
                        split_file_paths.append(output_filename)
                    print(f"Split PDF into {len(split_file_paths)} files based on ranges.")

                else:
                    raise ValueError(f"Unsupported split mode: {split_mode}")

                # Zip the created PDF files
                if split_file_paths:
                    with zipfile.ZipFile(output_zip_path, 'w') as zipf:
                        for file_path in split_file_paths:
                            zipf.write(file_path, os.path.basename(file_path))
                    print(f"Created ZIP file with split PDFs at: {output_zip_path}")
                else:
                    raise ValueError("No valid pages were selected for splitting.")

            return True
        finally:
            # Cancel the alarm
            signal.alarm(0)
            
    except TimeoutError as e:
        print(f"PDF split operation timed out: {str(e)}")
        return False
    except Exception as e:
        print(f"Error in split_pdf: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def compress_pdf(input_path, output_path, options=None):
    """
    Compress PDF using Ghostscript with different compression levels
    """
    try:
        import subprocess
        import json
        
        if options is None:
            options = {}
        
        compression_level = options.get('compression_level', 'ebook')
        grayscale = options.get('grayscale', False)
        
        print(f"Starting PDF compression: {input_path} -> {output_path}")
        print(f"Compression level: {compression_level}")
        print(f"Grayscale: {grayscale}")
        
        # Define Ghostscript parameters based on compression level
        compression_settings = {
            'screen': ['-dPDFSETTINGS=/screen', '-dColorImageResolution=72', '-dGrayImageResolution=72', '-dMonoImageResolution=72'],
            'ebook': ['-dPDFSETTINGS=/ebook', '-dColorImageResolution=150', '-dGrayImageResolution=150', '-dMonoImageResolution=150'],
            'printer': ['-dPDFSETTINGS=/printer', '-dColorImageResolution=300', '-dGrayImageResolution=300', '-dMonoImageResolution=300']
        }
        
        if compression_level not in compression_settings:
            print(f"Warning: Unknown compression level '{compression_level}', using 'ebook'")
            compression_level = 'ebook'
        
        # Base Ghostscript command
        gs_command = ['gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dNOPAUSE', '-dQUIET', '-dBATCH']
        
        # Add compression settings
        gs_command.extend(compression_settings[compression_level])
        
        # Add grayscale conversion if requested
        if grayscale:
            gs_command.extend(['-sProcessColorModel=DeviceGray', '-sColorConversionStrategy=Gray'])
        
        # Add input and output files
        gs_command.extend(['-sOutputFile=' + output_path, input_path])
        
        print(f"Ghostscript command: {' '.join(gs_command)}")
        
        # Execute Ghostscript with timeout
        result = subprocess.run(gs_command, capture_output=True, text=True, check=True, timeout=180)  # 3 minutes timeout
        
        print(f"Ghostscript stdout: {result.stdout}")
        if result.stderr:
            print(f"Ghostscript stderr: {result.stderr}")
        
        # Check if output file was created
        if os.path.exists(output_path):
            input_size = os.path.getsize(input_path)
            output_size = os.path.getsize(output_path)
            compression_ratio = (1 - output_size / input_size) * 100
            
            print(f"Compression successful!")
            print(f"Original size: {input_size:,} bytes")
            print(f"Compressed size: {output_size:,} bytes")
            print(f"Compression ratio: {compression_ratio:.1f}%")
            
            return True
        else:
            print(f"Error: Output file {output_path} was not created")
            return False
            
    except subprocess.CalledProcessError as e:
        print(f"Ghostscript error: {e}")
        print(f"Return code: {e.returncode}")
        print(f"stdout: {e.stdout}")
        print(f"stderr: {e.stderr}")
        return False
    except subprocess.TimeoutExpired as e:
        print(f"Ghostscript timeout: {e}")
        print(f"Command timed out after {e.timeout} seconds")
        return False
    except FileNotFoundError:
        print("Error: Ghostscript not found. Please install it:")
        print("  macOS: brew install ghostscript")
        print("  Ubuntu/Debian: sudo apt-get install ghostscript")
        print("  Windows: Download from https://www.ghostscript.com/")
        return False
    except Exception as e:
        print(f"Error in compress_pdf: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def protect_pdf(input_path, output_path, password):
    """
    Add password protection to PDF using PyPDF2
    """
    try:
        from PyPDF2 import PdfReader, PdfWriter
        
        print(f"Starting PDF protection: {input_path} -> {output_path}")
        print(f"Password length: {len(password)} characters")
        
        # Read the original PDF
        reader = PdfReader(input_path)
        writer = PdfWriter()
        
        # Add all pages to the writer
        for page in reader.pages:
            writer.add_page(page)
        
        # Encrypt with password
        writer.encrypt(password)
        
        # Write the protected PDF
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        print(f"PDF protection successful!")
        print(f"Original size: {os.path.getsize(input_path):,} bytes")
        print(f"Protected size: {os.path.getsize(output_path):,} bytes")
        
        return True
        
    except ImportError:
        print("Error: PyPDF2 not found. Please install it:")
        print("  pip install PyPDF2")
        return False
    except Exception as e:
        print(f"Error in protect_pdf: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def get_pdf_page_count(pdf_path):
    """
    Get the total number of pages in a PDF file
    """
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(pdf_path)
        return len(reader.pages)
    except Exception as e:
        print(f"Error getting page count: {str(e)}")
        return 0

def validate_page_operations(page_operations, total_pages):
    """
    Validate page operations against actual PDF page count
    Returns (is_valid, error_message)
    """
    errors = []
    
    # Validate page order
    page_order = page_operations.get('page_order', '')
    if page_order:
        try:
            page_numbers = [int(x.strip()) for x in page_order.split(',')]
            for page_num in page_numbers:
                if page_num < 1 or page_num > total_pages:
                    errors.append(f"Page {page_num} does not exist (PDF has {total_pages} pages)")
        except ValueError:
            errors.append("Invalid page order format. Use comma-separated numbers (e.g., 1,3,2,4)")
    
    # Validate rotate pages
    rotate_pages = page_operations.get('rotate_pages', '')
    if rotate_pages:
        try:
            for rotation in rotate_pages.split(','):
                if ':' in rotation:
                    page_num, angle = rotation.split(':')
                    page_num = int(page_num.strip())
                    angle = int(angle.strip())
                    
                    if page_num < 1 or page_num > total_pages:
                        errors.append(f"Page {page_num} does not exist for rotation (PDF has {total_pages} pages)")
                    
                    if angle not in [90, 180, 270]:
                        errors.append(f"Invalid rotation angle {angle}. Use 90, 180, or 270 degrees")
                else:
                    errors.append("Invalid rotation format. Use 'page:angle' (e.g., 2:90,3:180)")
        except ValueError:
            errors.append("Invalid rotation format. Use 'page:angle' (e.g., 2:90,3:180)")
    
    # Validate delete pages
    delete_pages = page_operations.get('delete_pages', '')
    if delete_pages:
        try:
            page_numbers = [int(x.strip()) for x in delete_pages.split(',')]
            for page_num in page_numbers:
                if page_num < 1 or page_num > total_pages:
                    errors.append(f"Page {page_num} does not exist for deletion (PDF has {total_pages} pages)")
        except ValueError:
            errors.append("Invalid delete pages format. Use comma-separated numbers (e.g., 4,5)")
    
    return len(errors) == 0, errors

def reorder_pages_pdf(input_path, output_path, page_operations):
    """
    Reorder, rotate, or delete pages in PDF using PyPDF2
    page_operations format: {
        "page_order": "1,3,2,4",  # Optional: reorder pages
        "rotate_pages": "2:90,3:180",  # Optional: rotate specific pages
        "delete_pages": "4,5"  # Optional: delete specific pages
    }
    """
    try:
        from PyPDF2 import PdfReader, PdfWriter
        
        print(f"Starting PDF page organization: {input_path} -> {output_path}")
        print(f"Page operations: {page_operations}")
        
        # Get total pages first
        total_pages = get_pdf_page_count(input_path)
        if total_pages == 0:
            return False
        
        print(f"Original PDF has {total_pages} pages")
        
        # Validate page operations
        is_valid, errors = validate_page_operations(page_operations, total_pages)
        if not is_valid:
            error_message = "; ".join(errors)
            print(f"Validation errors: {error_message}")
            raise ValueError(error_message)
        
        # Read the original PDF
        reader = PdfReader(input_path)
        
        # Parse page operations
        page_order = page_operations.get('page_order', '')
        rotate_pages = page_operations.get('rotate_pages', '')
        delete_pages = page_operations.get('delete_pages', '')
        
        # Create writer
        writer = PdfWriter()
        
        # Determine which pages to include and in what order
        pages_to_include = []
        
        if page_order:
            # Parse page order (e.g., "1,3,2,4")
            try:
                page_indices = [int(x.strip()) - 1 for x in page_order.split(',')]
                pages_to_include = page_indices
                print(f"Reordering pages: {page_order}")
            except ValueError:
                print("Invalid page order format, using all pages in original order")
                pages_to_include = list(range(total_pages))
        else:
            # Use all pages in original order
            pages_to_include = list(range(total_pages))
        
        # Parse pages to delete
        pages_to_delete = set()
        if delete_pages:
            try:
                delete_indices = [int(x.strip()) - 1 for x in delete_pages.split(',')]
                pages_to_delete = set(delete_indices)
                print(f"Deleting pages: {delete_pages}")
            except ValueError:
                print("Invalid delete pages format, skipping deletion")
        
        # Parse pages to rotate
        rotate_map = {}
        if rotate_pages:
            try:
                for rotation in rotate_pages.split(','):
                    if ':' in rotation:
                        page_num, angle = rotation.split(':')
                        page_index = int(page_num.strip()) - 1
                        angle_value = int(angle.strip())
                        rotate_map[page_index] = angle_value
                print(f"Rotating pages: {rotate_pages}")
            except ValueError:
                print("Invalid rotate pages format, skipping rotation")
        
        # Process pages
        final_pages = []
        for page_index in pages_to_include:
            if page_index in pages_to_delete:
                print(f"Skipping deleted page {page_index + 1}")
                continue
            
            if 0 <= page_index < total_pages:
                page = reader.pages[page_index]
                
                # Apply rotation if specified
                if page_index in rotate_map:
                    angle = rotate_map[page_index]
                    page.rotate(angle)
                    print(f"Rotated page {page_index + 1} by {angle} degrees")
                
                final_pages.append(page)
                print(f"Added page {page_index + 1}")
            else:
                print(f"Warning: Page {page_index + 1} does not exist, skipping")
        
        # Add all processed pages to writer
        for page in final_pages:
            writer.add_page(page)
        
        # Write the organized PDF
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        print(f"PDF page organization successful!")
        print(f"Original pages: {total_pages}")
        print(f"Final pages: {len(final_pages)}")
        print(f"Original size: {os.path.getsize(input_path):,} bytes")
        print(f"Organized size: {os.path.getsize(output_path):,} bytes")
        
        return True
        
    except ImportError:
        print("Error: PyPDF2 not found. Please install it:")
        print("  pip install PyPDF2")
        return False
    except ValueError as e:
        print(f"Validation error: {str(e)}")
        return False
    except Exception as e:
        print(f"Error in reorder_pages_pdf: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def ocr_pdf(input_path, output_path, options=None):
    """
    Perform OCR on PDF to make scanned documents searchable and selectable
    """
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image, ImageEnhance
        import fitz  # PyMuPDF
        import json
        import os
        
        # Set TESSDATA_PREFIX environment variable for Linux installation
        # In Docker container, tesseract is installed in the default location
        os.environ['TESSDATA_PREFIX'] = '/usr/share/tessdata'
        
        # Debug: Check if tesseract and language data are available
        import subprocess
        try:
            # Check tesseract version
            result = subprocess.run(['tesseract', '--version'], capture_output=True, text=True, timeout=10)
            print(f"Tesseract version: {result.stdout.split()[1] if result.stdout else 'Unknown'}")
            
            # Check if language data exists
            if os.path.exists('/usr/share/tessdata/eng.traineddata'):
                print("English language data found at /usr/share/tessdata/eng.traineddata")
            else:
                print("WARNING: English language data not found at /usr/share/tessdata/eng.traineddata")
                # List what's in the tessdata directory
                if os.path.exists('/usr/share/tessdata'):
                    files = os.listdir('/usr/share/tessdata')
                    print(f"Files in /usr/share/tessdata: {files}")
                else:
                    print("TESSDATA directory does not exist")
                    
        except Exception as e:
            print(f"Error checking tesseract installation: {e}")
        
        print(f"Starting OCR processing of {input_path}")
        
        if options is None:
            options = {}
        
        # Parse options if it's a JSON string
        if isinstance(options, str):
            try:
                options = json.loads(options)
            except json.JSONDecodeError as e:
                print(f"Error parsing options JSON: {e}")
                options = {}
        
        # Get OCR options
        language = options.get('language', 'eng')
        confidence = options.get('confidence', 0.7)
        preserve_layout = options.get('preserve_layout', True)
        extract_images = options.get('extract_images', True)
        
        print(f"OCR settings: language={language}, confidence={confidence}, preserve_layout={preserve_layout}")
        
        # Convert PDF pages to images with higher DPI for better accuracy
        print("Converting PDF pages to images...")
        images = convert_from_path(input_path, dpi=350)  # Good balance of quality and speed
        
        # Create new PDF with OCR text
        doc = fitz.open()
        
        for page_num, image in enumerate(images):
            print(f"Processing page {page_num + 1} with OCR...")
            
            # Simple image enhancement for better OCR accuracy
            enhanced_image = enhance_image_simple(image)
            
            # Perform OCR with optimized configuration
            ocr_config = f'--oem 3 --psm 6 -l {language}'
            
            # Check if language data exists and set tessdata directory accordingly
            tessdata_dir = '/usr/share/tessdata'
            if not os.path.exists(f'{tessdata_dir}/{language}.traineddata'):
                print(f"Warning: {language}.traineddata not found in {tessdata_dir}")
                # Try alternative locations
                alternative_dirs = ['/usr/local/share/tessdata', '/opt/tesseract/share/tessdata']
                for alt_dir in alternative_dirs:
                    if os.path.exists(f'{alt_dir}/{language}.traineddata'):
                        tessdata_dir = alt_dir
                        print(f"Found language data in {alt_dir}")
                        break
                else:
                    print("Warning: Language data not found in any standard location")
                    # Try without specifying tessdata directory
                    tessdata_dir = None
            
            if tessdata_dir:
                ocr_config += f' --tessdata-dir {tessdata_dir}'
            
            print(f"Using OCR config: {ocr_config}")
            
            # Get OCR data with positioning
            try:
                ocr_data = pytesseract.image_to_data(enhanced_image, config=ocr_config, output_type=pytesseract.Output.DICT)
            except Exception as ocr_error:
                print(f"OCR error with config '{ocr_config}': {ocr_error}")
                # Try with simpler config
                print("Trying with simpler OCR config...")
                ocr_data = pytesseract.image_to_data(enhanced_image, lang=language, output_type=pytesseract.Output.DICT)
            
            # Create new page with same dimensions as original
            page = doc.new_page(width=image.width, height=image.height)
            
            if preserve_layout:
                # Simple text positioning and filtering
                for i, conf in enumerate(ocr_data['conf']):
                    if conf > confidence * 100:  # Convert confidence to percentage
                        x = ocr_data['left'][i]
                        y = ocr_data['top'][i]
                        w = ocr_data['width'][i]
                        h = ocr_data['height'][i]
                        word = ocr_data['text'][i].strip()
                        
                        if word and h > 5:  # Basic filtering
                            # Calculate font size based on word height
                            font_size = max(8, min(14, h * 0.7))
                            page.insert_text((x, y + h), word, fontsize=font_size)
            
            # Add image if extract_images is enabled
            if extract_images:
                # Convert PIL image to fitz image using JPEG
                import io
                img_buffer = io.BytesIO()
                image.save(img_buffer, format='JPEG', quality=95)
                img_data = img_buffer.getvalue()
                page.insert_image(fitz.Rect(0, 0, image.width, image.height), stream=img_data)
            
            print(f"Completed OCR for page {page_num + 1}")
        
        # Save the OCR-processed PDF
        doc.save(output_path)
        doc.close()
        
        print(f"Successfully created OCR-processed PDF: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error during OCR processing: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def enhance_image_simple(image):
    """
    Simple image enhancement using only PIL for better OCR accuracy
    """
    try:
        # Convert to grayscale if not already
        if image.mode != 'L':
            image = image.convert('L')
        
        # Enhance contrast
        enhancer = ImageEnhance.Contrast(image)
        enhanced = enhancer.enhance(1.3)  # Increase contrast by 30%
        
        # Enhance sharpness slightly
        enhancer = ImageEnhance.Sharpness(enhanced)
        enhanced = enhancer.enhance(1.1)  # Slight sharpening
        
        return enhanced
        
    except Exception as e:
        print(f"Error in image enhancement: {e}")
        return image

def main():
    # Check for special commands that don't use argparse
    if len(sys.argv) > 1:
        command = sys.argv[1]
        
        if command == 'organize_pdf':
            input_path = sys.argv[2]
            output_path = sys.argv[3]
            
            # Read page operations from stdin
            page_operations = json.loads(sys.stdin.read())
            
            success = reorder_pages_pdf(input_path, output_path, page_operations)
            sys.exit(0 if success else 1)
        
        elif command == 'get_page_count':
            input_path = sys.argv[2]
            page_count = get_pdf_page_count(input_path)
            print(page_count)
            sys.exit(0 if page_count > 0 else 1)
    
    # Special handling for merge-pdf, which has a different argument structure
    if len(sys.argv) > 1 and sys.argv[1] == 'merge-pdf':
        if len(sys.argv) < 5: # script_name, merge-pdf, output_path, input_path1, ...
            print("Error: For merge-pdf, you must provide an output path and at least two input PDFs.")
            sys.exit(1)
        
        output_path = sys.argv[2]
        input_paths = sys.argv[3:]
        
        print(f"Python script received conversion_type: merge-pdf")
        print(f"Python script received output_path: {output_path}")
        print(f"Python script received input_paths: {', '.join(input_paths)}")
        
        success = merge_pdfs(output_path, input_paths)
        if success:
            print("Merge completed successfully")
            sys.exit(0)
        else:
            print("Merge failed")
            sys.exit(1)

    parser = argparse.ArgumentParser(description='Convert PDF to DOCX or Excel')
    parser.add_argument('conversion_type', choices=['pdf-to-word', 'pdf-to-excel', 'pdf-to-powerpoint', 'pdf-to-powerpoint-text', 'pdf-to-text', 'pdf-to-html', 'pdf-to-epub', 'pdf-to-rtf', 'pdf-to-svg', 'split-pdf', 'compress-pdf', 'protect-pdf', 'reorder-pages', 'ocr-pdf'],
                       help='Type of conversion to perform')
    parser.add_argument('pdf_path', help='Path to input PDF file')
    parser.add_argument('output_path', help='Path to output file')
    parser.add_argument('--page-selection', default='all', 
                       help='Page selection for Excel conversion (e.g., "1,3-5,7" or "all")')
    parser.add_argument('--options', default='{}',
                       help='JSON string of options for conversion')
    parser.add_argument('--password', default='',
                       help='Password for PDF protection')
    
    args = parser.parse_args()
    
    # Debug information
    print(f"Python script received conversion_type: {args.conversion_type}")
    print(f"Python script received pdf_path: {args.pdf_path}")
    print(f"Python script received output_path: {args.output_path}")
    print(f"Python script received page_selection: {args.page_selection}")
    print(f"Python script received options: {args.options}")
    print(f"Python script received password: {args.password}")
    print(f"Current working directory: {os.getcwd()}")
    print(f"File exists check result: {os.path.exists(args.pdf_path)}")
    
    # Parse options if provided
    options = {}
    if args.options and args.options != '{}':
        try:
            options = json.loads(args.options)
            print(f"Parsed options: {options}")
        except json.JSONDecodeError as e:
            print(f"Warning: Could not parse options JSON: {e}")
    
    # List files in the directory
    try:
        dir_path = os.path.dirname(args.pdf_path)
        print(f"Files in directory {dir_path}:")
        if os.path.exists(dir_path):
            for file in os.listdir(dir_path):
                print(f"  - {file}")
        else:
            print(f"Directory {dir_path} does not exist")
    except Exception as e:
        print(f"Error listing directory: {e}")
    
    # Check if input file exists
    if not os.path.exists(args.pdf_path):
        print(f"Error: Input file {args.pdf_path} does not exist")
        sys.exit(1)
    
    # Perform conversion based on type
    if args.conversion_type == 'pdf-to-word':
        success = pdf_to_word(args.pdf_path, args.output_path)
    elif args.conversion_type == 'pdf-to-excel':
        success = pdf_to_excel(args.pdf_path, args.output_path, args.page_selection)
    elif args.conversion_type == 'pdf-to-powerpoint':
        success = pdf_to_powerpoint(args.pdf_path, args.output_path)
    elif args.conversion_type == 'pdf-to-powerpoint-text':
        success = pdf_to_powerpoint_text(args.pdf_path, args.output_path)
    elif args.conversion_type == 'pdf-to-text':
        success = pdf_to_text(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'pdf-to-html':
        success = pdf_to_html(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'pdf-to-epub':
        success = pdf_to_epub(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'pdf-to-rtf':
        success = pdf_to_rtf(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'pdf-to-svg':
        success = pdf_to_svg(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'split-pdf':
        success = split_pdf(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'compress-pdf':
        success = compress_pdf(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'protect-pdf':
        success = protect_pdf(args.pdf_path, args.output_path, args.password)
    elif args.conversion_type == 'reorder-pages':
        success = reorder_pages_pdf(args.pdf_path, args.output_path, options)
    elif args.conversion_type == 'ocr-pdf':
        success = ocr_pdf(args.pdf_path, args.output_path, options)
    else:
        print(f"Error: Unknown conversion type {args.conversion_type}")
        sys.exit(1)
    
    if success:
        print("Conversion completed successfully")
        sys.exit(0)
    else:
        print("Conversion failed")
        sys.exit(1)

if __name__ == "__main__":
    main() 